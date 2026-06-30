from flask import Flask, request, jsonify, send_from_directory, send_file, session
from flask_cors import CORS
import ollama
import chromadb
from chromadb.config import Settings
import os
import pandas as pd
import PyPDF2
import pdfplumber
from docx import Document
from pptx import Presentation
from sentence_transformers import SentenceTransformer
import shutil
from datetime import datetime, timedelta
import uuid
import io
import json
import sqlite3
import hashlib
from functools import wraps
from werkzeug.utils import secure_filename
import secrets

app = Flask(__name__)
CORS(app, supports_credentials=True, origins='*', allow_headers=['Content-Type', 'Authorization'], methods=['GET', 'POST', 'PUT', 'DELETE', 'OPTIONS'])

# Session configuration - use cookie-based sessions
app.config['SECRET_KEY'] = 'your-secret-key-change-this-in-production'
app.config['SESSION_COOKIE_HTTPONLY'] = False
app.config['SESSION_COOKIE_SAMESITE'] = 'Lax'
app.config['PERMANENT_SESSION_LIFETIME'] = timedelta(hours=24)

# Configuration
UPLOAD_FOLDER = 'uploads'
VECTOR_DB_PATH = 'chroma_db'
DATABASE_PATH = 'rbac.db'
ALLOWED_EXTENSIONS = {'xlsx', 'xls', 'pdf', 'docx'}
COMPANIES = ['A', 'B', 'C', 'D', 'HQ', 'SP']

# Create directories
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(VECTOR_DB_PATH, exist_ok=True)

# Initialize ChromaDB
chroma_client = chromadb.PersistentClient(path=VECTOR_DB_PATH)
collection = chroma_client.get_or_create_collection(name="documents")

# Initialize embedding model
embedding_model = SentenceTransformer('all-MiniLM-L6-v2')

# Conversation history storage (in-memory for simplicity)
conversation_history = {}  # {session_id: [{'role': 'user'|'assistant', 'content': str}]}

# Database initialization
def init_db():
    conn = sqlite3.connect(DATABASE_PATH)
    cursor = conn.cursor()
    
    # Users table
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT UNIQUE NOT NULL,
            password_hash TEXT NOT NULL,
            full_name TEXT NOT NULL,
            role TEXT NOT NULL CHECK (role IN ('admin', 'user')),
            companies TEXT NOT NULL,
            status TEXT NOT NULL DEFAULT 'active' CHECK (status IN ('active', 'inactive')),
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    
    # Files table
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS files (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            filename TEXT NOT NULL,
            original_filename TEXT NOT NULL,
            company TEXT NOT NULL CHECK (company IN ('A', 'B', 'C', 'D', 'HQ', 'SP')),
            file_size INTEGER NOT NULL,
            upload_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            uploaded_by INTEGER NOT NULL,
            status TEXT NOT NULL DEFAULT 'active' CHECK (status IN ('active', 'deleted')),
            FOREIGN KEY (uploaded_by) REFERENCES users (id)
        )
    ''')
    
    # Audit logs table
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS audit_logs (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER,
            action TEXT NOT NULL,
            resource_type TEXT NOT NULL,
            resource_id TEXT,
            details TEXT,
            ip_address TEXT,
            timestamp TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (user_id) REFERENCES users (id)
        )
    ''')
    
    # AI Training table
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS ai_training (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            company TEXT CHECK (company IN ('A', 'B', 'C', 'D', 'HQ', 'SP', 'ALL')),
            instruction TEXT NOT NULL,
            category TEXT NOT NULL,
            priority INTEGER DEFAULT 5,
            created_by INTEGER NOT NULL,
            created_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            updated_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            enabled BOOLEAN DEFAULT 1,
            FOREIGN KEY (created_by) REFERENCES users (id)
        )
    ''')
    
    # Excel Schema Metadata table
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS excel_schema (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            file_id INTEGER NOT NULL,
            column_name TEXT NOT NULL,
            data_type TEXT NOT NULL,
            detected_meaning TEXT,
            synonyms TEXT,
            sample_values TEXT,
            is_key_column BOOLEAN DEFAULT 0,
            is_date_column BOOLEAN DEFAULT 0,
            is_numeric_column BOOLEAN DEFAULT 0,
            created_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (file_id) REFERENCES files (id)
        )
    ''')
    
    # Create default admin user if not exists
    cursor.execute('SELECT * FROM users WHERE username = ?', ('admin',))
    if not cursor.fetchone():
        password_hash = hashlib.sha256('admin123'.encode()).hexdigest()
        cursor.execute('''
            INSERT INTO users (username, password_hash, full_name, role, companies)
            VALUES (?, ?, ?, ?, ?)
        ''', ('admin', password_hash, 'System Administrator', 'admin', 'A,B,C,D,HQ,SP'))
    
    conn.commit()
    conn.close()

# Initialize database on startup
init_db()

# Database helper functions
def get_db_connection():
    conn = sqlite3.connect(DATABASE_PATH)
    conn.row_factory = sqlite3.Row
    return conn

# Ensure ai_training table exists (for existing databases)
def ensure_training_table():
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS ai_training (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            company TEXT NOT NULL DEFAULT 'ALL',
            instruction TEXT NOT NULL,
            category TEXT NOT NULL,
            priority INTEGER DEFAULT 5,
            created_by INTEGER NOT NULL,
            created_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            updated_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            enabled BOOLEAN DEFAULT 1,
            FOREIGN KEY (created_by) REFERENCES users (id)
        )
    ''')
    conn.commit()
    conn.close()

ensure_training_table()

def hash_password(password):
    return hashlib.sha256(password.encode()).hexdigest()

def verify_password(password, password_hash):
    return hashlib.sha256(password.encode()).hexdigest() == password_hash

def log_audit(user_id, action, resource_type, resource_id=None, details=None, ip_address=None):
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute('''
        INSERT INTO audit_logs (user_id, action, resource_type, resource_id, details, ip_address)
        VALUES (?, ?, ?, ?, ?, ?)
    ''', (user_id, action, resource_type, resource_id, details, ip_address))
    conn.commit()
    conn.close()

def analyze_and_store_excel_schema(file_id, table_data, filepath):
    """Analyze Excel schema and store metadata in database"""
    if not table_data or 'columns' not in table_data or 'data' not in table_data:
        return
    
    columns = table_data['columns']
    data = table_data['data']
    
    conn = get_db_connection()
    cursor = conn.cursor()
    
    try:
        for col_name in columns:
            col_index = columns.index(col_name)
            
            # Detect data type from sample values
            sample_values = []
            data_types = set()
            
            for row in data[:min(100, len(data))]:  # Sample first 100 rows
                if col_index < len(row):
                    value = row[col_index]
                    # Convert datetime objects to ISO format strings for JSON serialization
                    if hasattr(value, 'isoformat'):
                        value = value.isoformat()
                    else:
                        value = str(value)
                    sample_values.append(value)
                    
                    # Detect data type
                    if value is None or str(value).strip() == '':
                        continue
                    elif isinstance(value, (int, float)):
                        data_types.add('numeric')
                    elif isinstance(value, str):
                        # Check for date patterns
                        if any(pattern in str(value).lower() for pattern in ['date', 'time', '/', '-']):
                            data_types.add('date')
                        else:
                            data_types.add('text')
            
            # Determine primary data type
            primary_type = 'text'
            if 'numeric' in data_types:
                primary_type = 'numeric'
            elif 'date' in data_types:
                primary_type = 'date'
            
            # Detect column meaning using heuristics
            detected_meaning = detect_column_meaning(col_name, sample_values)
            
            # Generate synonyms
            synonyms = generate_column_synonyms(col_name, detected_meaning)
            
            # Determine if key column
            is_key = is_likely_key_column(col_name, sample_values)
            
            # Store schema metadata
            cursor.execute('''
                INSERT INTO excel_schema 
                (file_id, column_name, data_type, detected_meaning, synonyms, sample_values, is_key_column, is_date_column, is_numeric_column)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
            ''', (
                file_id,
                col_name,
                primary_type,
                detected_meaning,
                json.dumps(synonyms),
                json.dumps(sample_values[:20]),  # Store first 20 samples
                1 if is_key else 0,
                1 if primary_type == 'date' else 0,
                1 if primary_type == 'numeric' else 0
            ))
        
        conn.commit()
        print(f"Excel schema analyzed and stored for file_id: {file_id}")
    except Exception as e:
        print(f"Error analyzing Excel schema: {e}")
        conn.rollback()
    finally:
        conn.close()

def detect_column_meaning(col_name, sample_values):
    """Detect the meaning of a column based on name and sample values"""
    col_lower = col_name.lower()
    
    # Common column patterns - updated for military personnel database
    patterns = {
        'name': ['name', 'person', 'employee', 'soldier', 'individual', 'wife name'],
        'qualification': ['qual', 'education', 'degree', 'btech', 'mba', 'phd', 'b.e', 'm.e', 'couse', 'course'],
        'sports': ['sport', 'game', 'athletic', 'football', 'cricket', 'basketball', 'tennis', 'sports event'],
        'experience': ['exp', 'service', 'years', 'tenure', 'duration', 'svc', 'ere', 'appt duty'],
        'department': ['dept', 'unit', 'company', 'branch', 'division', 'coy'],
        'rank': ['rank', 'designation', 'position', 'grade', 'level'],
        'medical': ['medical', 'health', 'shape', 'category', 'fitness', 'med cat', 'blood gp'],
        'trade': ['trade', 'skill', 'specialization', 'technical', 'signals', 'it', 'electronics'],
        'location': ['location', 'place', 'station', 'post', 'area', 'vill', 'teh', 'distt', 'state', 'pin code'],
        'id': ['id', 'code', 'number', 'roll', 'reg', 'emp', 'army no', 'i card number', 'uid no', 'pan number', 'sr no'],
        'date': ['date', 'time', 'dob', 'joining', 'birth', 'doe', 'dor', 'dom', 'from', 'to', 'dob wife', 'dob ch1'],
        'age': ['age', 'years_old', 'yrs'],
        'contact': ['contact', 'phone', 'mobile', 'email', 'address', 'mob no', 'e mail id'],
        'family': ['wife', 'children', 'marital', 'cast', 'achievement', 'punishment'],
        'financial': ['acct', 'account', 'bank', 'single acct', 'jt account']
    }
    
    for meaning, keywords in patterns.items():
        if any(keyword in col_lower for keyword in keywords):
            return meaning.capitalize()
    
    # Check sample values for patterns
    if sample_values:
        # Check for qualification patterns
        if any('btech' in v.lower() or 'b.tech' in v.lower() or 'b.e' in v.lower() for v in sample_values):
            return 'Qualification'
        # Check for sports patterns
        if any('football' in v.lower() or 'cricket' in v.lower() or 'basketball' in v.lower() for v in sample_values):
            return 'Sports'
        # Check for medical categories
        if any('shape' in v.lower() or 'medical' in v.lower() for v in sample_values):
            return 'Medical'
    
    return col_name.capitalize()

def generate_column_synonyms(col_name, detected_meaning):
    """Generate synonyms for a column name - updated for military personnel database"""
    synonyms = [col_name]
    col_lower = col_name.lower()
    
    # Common synonym mappings - updated for military columns
    synonym_map = {
        'name': ['full name', 'person name', 'employee name', 'soldier name', 'individual name'],
        'qualification': ['education', 'degree', 'academic', 'qual', 'course', 'courses completed'],
        'sports': ['games', 'athletics', 'sports played', 'activities', 'sports event'],
        'experience': ['service years', 'tenure', 'years of service', 'work experience', 'svc yrs', 'ere', 'extra regimental duty'],
        'department': ['dept', 'unit', 'company', 'branch', 'division', 'coy', 'company'],
        'rank': ['designation', 'position', 'grade', 'level'],
        'medical': ['health', 'medical category', 'shape', 'fitness', 'med cat', 'blood group'],
        'trade': ['skill', 'specialization', 'technical', 'trade'],
        'service': ['service years', 'tenure', 'years served'],
        'id': ['code', 'number', 'roll number', 'employee id', 'army number', 'army no', 'identity card', 'i card', 'uid', 'aadhar', 'pan', 'serial number', 'sr no'],
        'date': ['date of birth', 'birth date', 'joining date', 'enrollment date', 'retirement date', 'marriage date', 'dob', 'doe', 'dor', 'dom'],
        'contact': ['phone', 'mobile', 'email', 'address', 'mob no', 'mobile number', 'email id'],
        'family': ['wife', 'spouse', 'children', 'kids', 'marital status', 'caste', 'achievements', 'awards', 'punishment record'],
        'financial': ['account', 'bank account', 'bank details', 'single account', 'joint account', 'acct no'],
        'location': ['village', 'post office', 'tehsil', 'district', 'state', 'pincode', 'address', 'residence']
    }
    
    # Specific column name mappings
    specific_mappings = {
        'army no': ['army number', 'service number', 'military id'],
        'i card number': ['identity card', 'id card', 'icard'],
        'uid no': ['aadhar', 'aadhaar', 'uid'],
        'pan number': ['pan', 'pan card'],
        'blood gp': ['blood group', 'blood type'],
        'med cat': ['medical category', 'med category', 'shape'],
        'coy': ['company', 'unit', 'coy'],
        'dob': ['date of birth', 'birth date'],
        'doe': ['date of enrollment', 'enrollment date', 'joining date'],
        'dor': ['date of retirement', 'retirement date'],
        'dom': ['date of marriage', 'marriage date'],
        'mob no': ['mobile number', 'phone number', 'contact number'],
        'e mail id': ['email', 'email address', 'mail id'],
        'single acct no': ['single account', 'personal account', 'individual account'],
        'jt account no': ['joint account', 'joint bank account'],
        'couse': ['course', 'qualification', 'training'],
        'ere': ['extra regimental duty', 'additional duty', 'ere duty'],
        'appt duty': ['appointment', 'current duty', 'present posting', 'current appointment'],
        'sports event': ['sports', 'games', 'athletics', 'sports played'],
        'achievement': ['awards', 'medals', 'honors', 'achievements'],
        'punishment': ['disciplinary action', 'penalty', 'punishment record'],
        'marital status': ['married', 'unmarried', 'marriage status'],
        'vill': ['village', 'native place'],
        'post': ['post office', 'po'],
        'teh': ['tehsil', 'taluka'],
        'distt': ['district'],
        'pin code': ['pincode', 'postal code', 'zip'],
        'wife name': ['spouse name', 'wife'],
        'dob wife': ['wife dob', 'wife date of birth', 'spouse dob'],
        'children1': ['first child', 'child 1', 'eldest child'],
        'dob ch1': ['first child dob', 'child 1 dob'],
        'children 2': ['second child', 'child 2']
    }
    
    # Add specific mappings
    for key, syns in specific_mappings.items():
        if key in col_lower:
            synonyms.extend(syns)
    
    # Add general synonym mappings
    for key, syns in synonym_map.items():
        if key in col_lower or key in detected_meaning.lower():
            synonyms.extend(syns)
    
    return list(set(synonyms))

def is_likely_key_column(col_name, sample_values):
    """Determine if a column is likely a key/identifier column"""
    col_lower = col_name.lower()
    
    # Check for ID patterns
    if any(keyword in col_lower for keyword in ['id', 'code', 'number', 'roll', 'reg', 'emp']):
        return True
    
    # Check if values are unique
    if sample_values and len(set(sample_values)) == len(sample_values):
        return True
    
    return False

def get_column_mapping_for_query(query, user_companies):
    """Get intelligent column mapping based on query and stored schema metadata"""
    conn = get_db_connection()
    cursor = conn.cursor()
    
    try:
        # Get all schema metadata for user's accessible companies
        placeholders = ','.join(['?'] * len(user_companies))
        cursor.execute(f'''
            SELECT es.column_name, es.detected_meaning, es.synonyms, f.company
            FROM excel_schema es
            JOIN files f ON es.file_id = f.id
            WHERE f.company IN ({placeholders})
        ''', user_companies)
        
        schema_data = cursor.fetchall()
        
        # Build mapping from query terms to actual column names
        column_mapping = {}
        query_lower = query.lower()
        
        for row in schema_data:
            col_name = row['column_name']
            detected_meaning = row['detected_meaning']
            synonyms = json.loads(row['synonyms']) if row['synonyms'] else []
            
            # Check if any synonym appears in query
            for synonym in synonyms + [col_name, detected_meaning]:
                if synonym.lower() in query_lower:
                    column_mapping[synonym] = col_name
                    column_mapping[detected_meaning.lower()] = col_name
        
        return column_mapping
    except Exception as e:
        print(f"Error getting column mapping: {e}")
        return {}
    finally:
        conn.close()

# Authentication decorators
def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        print(f"Session check: {dict(session)}")
        if 'user_id' not in session:
            return jsonify({'error': 'Authentication required'}), 401
        return f(*args, **kwargs)
    return decorated_function

def admin_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'user_id' not in session:
            return jsonify({'error': 'Authentication required'}), 401
        if session.get('role') != 'admin':
            return jsonify({'error': 'Admin access required'}), 403
        return f(*args, **kwargs)
    return decorated_function

def get_user_companies(user_id):
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute('SELECT companies FROM users WHERE id = ?', (user_id,))
    result = cursor.fetchone()
    conn.close()
    if result:
        return result['companies'].split(',')
    return []

def check_company_access(user_id, company):
    if session.get('role') == 'admin':
        return True
    user_companies = get_user_companies(user_id)
    return company in user_companies

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pdf(file_path):
    text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text += page.extract_text() + "\n"
    except:
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            print(f"Error reading PDF: {e}")
    return text

def extract_text_from_word(file_path):
    text = ""
    try:
        doc = Document(file_path)
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
    except Exception as e:
        print(f"Error reading Word document: {e}")
    return text

def extract_text_from_excel(file_path):
    text = ""
    table_data = None
    try:
        df = pd.read_excel(file_path)
        text = df.to_string()
        
        # Convert datetime objects to strings for JSON serialization
        def convert_value(val):
            if hasattr(val, 'isoformat'):
                return val.isoformat()
            elif pd.isna(val):
                return ""
            else:
                return str(val)
        
        # Store as JSON for table rendering with datetime conversion
        table_data = {
            'columns': df.columns.tolist(),
            'data': [[convert_value(val) for val in row] for row in df.values.tolist()],
            'index': df.index.tolist()
        }
    except Exception as e:
        print(f"Error reading Excel file: {e}")
    return text, table_data

def chunk_text(text, chunk_size=500, overlap=50):
    """Improved chunking with overlap for better context preservation"""
    chunks = []
    words = text.split()
    
    for i in range(0, len(words), chunk_size - overlap):
        chunk = ' '.join(words[i:i+chunk_size])
        if chunk.strip():  # Only add non-empty chunks
            chunks.append(chunk)
    
    return chunks

def generate_excel_file(data, filename):
    """Generate Excel file from data"""
    output = io.BytesIO()
    
    # Parse data - could be JSON or structured text
    try:
        # Try to parse as JSON first
        if isinstance(data, str):
            parsed_data = json.loads(data)
        else:
            parsed_data = data
        
        # Create DataFrame
        if isinstance(parsed_data, list):
            df = pd.DataFrame(parsed_data)
        elif isinstance(parsed_data, dict):
            df = pd.DataFrame([parsed_data])
        else:
            df = pd.DataFrame({'Content': [str(data)]})
    except:
        # If not JSON, create simple DataFrame
        df = pd.DataFrame({'Content': [str(data)]})
    
    # Save to Excel
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        df.to_excel(writer, index=False, sheet_name='Data')
    
    output.seek(0)
    return output

def generate_word_file(content, filename):
    """Generate Word file from content"""
    doc = Document()
    
    # Split content into paragraphs
    paragraphs = content.split('\n')
    for para in paragraphs:
        if para.strip():
            doc.add_paragraph(para.strip())
    
    output = io.BytesIO()
    doc.save(output)
    output.seek(0)
    return output

def generate_ppt_file(content, filename):
    """Generate PowerPoint file from content"""
    prs = Presentation()
    
    # Add title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = "AI Generated Report"
    subtitle.text = datetime.now().strftime("%Y-%m-%d %H:%M")
    
    # Add content slides
    paragraphs = content.split('\n')
    current_slide = None
    current_text = ""
    
    for para in paragraphs:
        if para.strip():
            current_text += para.strip() + "\n"
            if len(current_text) > 300:  # Create new slide every ~300 chars
                if current_slide is None:
                    current_slide = prs.slides.add_slide(prs.slide_layouts[1])
                else:
                    current_slide = prs.slides.add_slide(prs.slide_layouts[1])
                
                # Add text to slide
                text_box = current_slide.placeholders[1]
                text_box.text = current_text.strip()
                current_text = ""
    
    # Add remaining text
    if current_text.strip():
        current_slide = prs.slides.add_slide(prs.slide_layouts[1])
        text_box = current_slide.placeholders[1]
        text_box.text = current_text.strip()
    
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

@app.route('/')
def serve_index():
    return send_from_directory('.', 'index.html')

# Authentication endpoints
@app.route('/api/auth/login', methods=['POST'])
def login():
    data = request.json
    username = data.get('username')
    password = data.get('password')
    
    if not username or not password:
        return jsonify({'error': 'Username and password required'}), 400
    
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute('SELECT * FROM users WHERE username = ? AND status = ?', (username, 'active'))
    user = cursor.fetchone()
    conn.close()
    
    if not user or not verify_password(password, user['password_hash']):
        log_audit(None, 'login_failed', 'authentication', details=f'Failed login attempt for username: {username}', ip_address=request.remote_addr)
        return jsonify({'error': 'Invalid credentials'}), 401
    
    session['user_id'] = user['id']
    session['username'] = user['username']
    session['full_name'] = user['full_name']
    session['role'] = user['role']
    session['companies'] = user['companies'].split(',')
    
    print(f"Session set: user_id={session.get('user_id')}, role={session.get('role')}")
    
    log_audit(user['id'], 'login', 'authentication', ip_address=request.remote_addr)
    
    return jsonify({
        'message': 'Login successful',
        'user': {
            'id': user['id'],
            'username': user['username'],
            'full_name': user['full_name'],
            'role': user['role'],
            'companies': user['companies'].split(',')
        }
    }), 200

@app.route('/api/auth/logout', methods=['POST'])
@login_required
def logout():
    user_id = session.get('user_id')
    log_audit(user_id, 'logout', 'authentication', ip_address=request.remote_addr)
    session.clear()
    return jsonify({'message': 'Logout successful'}), 200

@app.route('/api/auth/me', methods=['GET'])
@login_required
def get_current_user():
    return jsonify({
        'user': {
            'id': session.get('user_id'),
            'username': session.get('username'),
            'full_name': session.get('full_name'),
            'role': session.get('role'),
            'companies': session.get('companies', [])
        }
    }), 200

@app.route('/api/upload', methods=['POST'])
@login_required
@admin_required
def upload_file():
    print(f"Upload request received. Session: {session}")
    print(f"Files in request: {list(request.files.keys())}")
    print(f"Form data: {list(request.form.keys())}")
    
    if 'file' not in request.files:
        print("Error: No file part in request")
        return jsonify({'error': 'No file part'}), 400
    
    file = request.files['file']
    if file.filename == '':
        print("Error: No selected file")
        return jsonify({'error': 'No selected file'}), 400
    
    company = request.form.get('company')
    print(f"Company from form: {company}")
    if not company or company not in COMPANIES:
        print(f"Error: Invalid company {company}")
        return jsonify({'error': 'Valid company required (A, B, C, D, HQ, SP)'}), 400
    
    if file and allowed_file(file.filename):
        try:
            filename = f"{uuid.uuid4()}_{file.filename}"
            filepath = os.path.join(UPLOAD_FOLDER, filename)
            file.save(filepath)
            print(f"File saved to: {filepath}")
            
            # Get file size
            file_size = os.path.getsize(filepath)
            
            # Extract text based on file type
            file_ext = filename.rsplit('.', 1)[1].lower()
            table_data = None
            if file_ext in ['xlsx', 'xls']:
                text, table_data = extract_text_from_excel(filepath)
            elif file_ext == 'pdf':
                text = extract_text_from_pdf(filepath)
            elif file_ext == 'docx':
                text = extract_text_from_word(filepath)
            
            print(f"Text extracted, length: {len(text)}")
            
            # Chunk the text
            chunks = chunk_text(text)
            print(f"Text chunked into {len(chunks)} chunks")
            
            # Generate embeddings and store in ChromaDB with company metadata
            embeddings = embedding_model.encode(chunks).tolist()
            ids = [f"{filename}_{i}" for i in range(len(chunks))]
            
            # Convert table_data to JSON string for ChromaDB metadata (only for first chunk)
            table_data_json = json.dumps(table_data) if table_data else ""
            
            metadatas = []
            for i in range(len(chunks)):
                metadata = {
                    "source": filename, 
                    "chunk": i, 
                    "is_excel": file_ext in ['xlsx', 'xls'], 
                    "company": company,
                    "uploaded_by": session.get('user_id')
                }
                # Only add table_data for first chunk if it exists
                if i == 0 and table_data_json:
                    metadata["table_data"] = table_data_json
                metadatas.append(metadata)
            
            collection.add(
                embeddings=embeddings,
                documents=chunks,
                ids=ids,
                metadatas=metadatas
            )
            print("Added to ChromaDB")
            
            # Store file metadata in database
            conn = get_db_connection()
            cursor = conn.cursor()
            cursor.execute('''
                INSERT INTO files (filename, original_filename, company, file_size, uploaded_by)
                VALUES (?, ?, ?, ?, ?)
            ''', (filename, file.filename, company, file_size, session.get('user_id')))
            file_id = cursor.lastrowid
            conn.commit()
            print(f"File metadata stored in database with ID: {file_id}")
            
            # Analyze Excel schema and store metadata after file_id is created
            if file_ext in ['xlsx', 'xls'] and table_data:
                analyze_and_store_excel_schema(file_id, table_data, filepath)
            
            conn.close()
            
            log_audit(session.get('user_id'), 'upload', 'file', str(file_id), 
                      details=f'Uploaded {file.filename} for company {company}', ip_address=request.remote_addr)
            
            return jsonify({
                'message': 'File uploaded and processed successfully',
                'filename': filename,
                'company': company,
                'chunks': len(chunks)
            }), 200
        except Exception as e:
            print(f"Error during upload: {str(e)}")
            import traceback
            traceback.print_exc()
            return jsonify({'error': f'Upload failed: {str(e)}'}), 500
    
    return jsonify({'error': 'Invalid file type'}), 400

@app.route('/api/chat-public', methods=['POST'])
def chat_public():
    """Public chat endpoint for Prachand Paanch page - no Flask auth required"""
    data = request.json
    query = data.get('message', '')
    format_type = data.get('format', 'text')
    companies = data.get('companies', [])  # Companies passed from frontend
    
    if not query:
        return jsonify({'error': 'No message provided'}), 400
    
    # Initialize conversation history (simplified for public endpoint)
    session_id = str(uuid.uuid4())
    
    # Retrieve relevant documents with company filtering
    query_embedding = embedding_model.encode([query]).tolist()
    
    # Filter by provided companies
    if companies and len(companies) > 0:
        results = collection.query(
            query_embeddings=query_embedding,
            n_results=10,
            where={"company": {"$in": companies}}
        )
    else:
        # If no companies specified, return all (or restrict based on your security needs)
        results = collection.query(
            query_embeddings=query_embedding,
            n_results=10
        )
    
    # Build context from retrieved documents
    context = ""
    table_data = None
    is_excel_source = False
    source_info = []
    
    if results['documents'] and results['documents'][0]:
        for i, (doc, meta) in enumerate(zip(results['documents'][0], results['metadatas'][0])):
            source = meta.get('source', 'Unknown')
            company = meta.get('company', 'Unknown')
            chunk_num = meta.get('chunk', i)
            
            source_info.append(f"[Source: {source}, Company: {company}, Chunk: {chunk_num}]")
            
            if meta.get('is_excel', False):
                is_excel_source = True
                if meta.get('table_data'):
                    try:
                        table_data = json.loads(meta['table_data'])
                    except:
                        table_data = None
            
            context += f"\n--- {source_info[-1]} ---\n{doc}\n"
        
        if source_info:
            context = f"Retrieved from {len(source_info)} document chunks:\n" + context
    
    # Use the same system prompt as the authenticated endpoint
    system_prompt = """You are an intelligent military personnel data analysis assistant. Your job is to provide accurate, well-formatted, professional responses based on user queries about military personnel records.

CRITICAL DATA ACCURACY RULES:

A. STRICT QUALIFICATION/COURSE MATCHING
- When user asks for specific qualification/course (e.g., "B.Tech", "Technical", "Course"), return ONLY personnel explicitly matching that qualification in the "couse/course" field
- Qualifications must be EXACTLY present in the course field
- Example: "B.Tech personnel" → Only include records where "B.Tech" is explicitly present in course field
- Courses like "B.Tech, M.Tech" → Include (B.Tech present)
- Courses like "MBA, B.Tech" → Include (B.Tech present)
- Courses like "MBA, PhD" → Exclude (B.Tech absent)
- NEVER infer or assume related qualifications
- If qualification-specific data not available: "[Qualification]-specific data is not available. Unable to determine [qualification] holders."
- For "technical qualification" queries: Only include B.Tech, M.Tech, B.E., M.E., or similar technical degrees. EXCLUDE MBA, PhD, Arts, Commerce, etc.

B. STRICT SPORTS MATCHING
- When user asks for specific sport (e.g., "Football", "Cricket", "Basketball"), return ONLY personnel where that sport exists in the "sports event" field
- Example: "Football players" → Only include records where "Football" is explicitly present in sports event column
- Sports like "Cricket, Football" → Include (Football present)
- Sports like "Football, Basketball" → Include (Football present)
- Sports like "Volleyball, Tennis" → Exclude (Football absent)
- NEVER infer suitability from other sports
- If sport-specific data not available: "[Sport]-specific data is not available. Unable to determine [sport] suitability."

C. STRICT FILTERING
- Apply ALL filters specified in the query
- If user asks for "Company A personnel" (coy field), ONLY return Company A records
- If user asks for "personnel with X rank", ONLY return records where X is present in rank field
- If user asks for "personnel from X village", ONLY return records where X is present in vill field
- NEVER include records that don't match specified criteria
- If no records match: "No records found matching the specified criteria."

D. NO ASSUMPTIONS
- NEVER assume or infer information not explicitly stated in the documents
- If information is missing, state it clearly: "Information not available in the documents"
- NEVER make up or hallucinate data
- If uncertain: "Based on available documents, [information] is not specified"

E. HALLUCINATION PREVENTION
- ONLY use information from the provided document context
- NEVER add external knowledge or assumptions
- If context doesn't contain the answer: "The provided documents do not contain information about [topic]"
- NEVER fabricate data, statistics, or details

F. QUERY VERIFICATION
- Before responding, verify that your answer is supported by the document context
- If query asks about specific data points and they're not in context: "The documents do not contain [specific data]"
- If query is ambiguous: ask for clarification or provide the most relevant available information

G. COLUMN SEPARATION
- For tabular/structured responses, ensure relevant fields are in separate columns
- NEVER combine different data types in a single column
- Format: | Army No | Name | Rank | Company | Qualification | Sports | ...

H. QUERY INTENT RULES
- Evaluate query intent in this order:
  1. Qualification/course-specific queries (e.g., "B.Tech personnel") → Apply STRICT QUALIFICATION MATCHING
  2. Sport-specific queries (e.g., "Football players") → Apply STRICT SPORTS MATCHING
  3. Company-specific queries (e.g., "Company A personnel") → Apply STRICT FILTERING
  4. Rank-specific queries (e.g., "Lieutenant personnel") → Apply STRICT FILTERING
  5. Location-specific queries (e.g., "personnel from X village") → Apply STRICT FILTERING
  6. General queries → Provide comprehensive information from context

I. DATE FIELD HANDLING
- When querying dates (dob, doe, dor, dom), use the actual date values from the respective columns
- DOB = Date of Birth
- DOE = Date of Enrollment
- DOR = Date of Retirement
- DOM = Date of Marriage
- For date range queries, compare against the appropriate date field

J. FAMILY INFORMATION HANDLING
- Wife information is in "wife name" and "dob wife" fields
- Children information is in "children1", "dob ch1", "children 2" fields
- Marital status is in "marital status" field
- Caste information is in "cast" field

ABSOLUTE FORMAT RULES:
- If user asks for "table", "tabular format", "tabular", "in a table": Provide ONLY a Markdown table. Nothing else.
- If user asks for "bullet points", "list", "points": Provide ONLY bullet points in "- Field: Value" format. Nothing else.
- If no format specified: Provide ONLY bullet points in "- Field: Value" format. Nothing else.
- NEVER provide multiple formats in one response
- NEVER add "Here is the information", "Let me know if you need anything else", or any other conversational text
- NEVER mix bullet points and tables
- Your response must contain ONLY the formatted data, nothing more

Example correct responses:
Table format:
| Army No | Name | Rank | Company | Qualification |
| --- | --- | --- | --- | --- |
| 12345 | John Doe | Lieutenant | A | B.Tech |

Bullet point format:
- Army No: 12345
- Name: John Doe
- Rank: Lieutenant
- Company: A
- Qualification: B.Tech

IMPORTANT SECURITY INSTRUCTIONS:
- Only provide information from the document context provided
- Do not make up or hallucinate information
- If the context is empty or doesn't contain relevant information, say so clearly
"""
    
    query_lower = query.lower()
    format_instruction = ""
    
    if any(keyword in query_lower for keyword in ['table', 'tabular format', 'tabular', 'in a table']):
        format_instruction = "IMPORTANT: Provide ONLY a clean Markdown table. No introductory text, no explanatory text, no concluding remarks. Just the table with proper headers and alignment."
    elif any(keyword in query_lower for keyword in ['bullet points', 'list', 'points', 'bullet']):
        format_instruction = "IMPORTANT: Format using bullet points with '- Field: Value' format for each data point. No introductory or explanatory text."
    else:
        format_instruction = "IMPORTANT: Format using bullet points with '- Field: Value' format for each data point. No introductory or explanatory text."
    
    prompt = f"""{system_prompt}

Document Context (from uploaded files):
{context}

Current Question: {query}

{format_instruction}

If the context doesn't contain relevant information, say so clearly."""
    
    messages = [{'role': 'system', 'content': system_prompt}, {'role': 'user', 'content': prompt}]
    
    try:
        response = ollama.chat(model='llama3', messages=messages)
        response_text = response['message']['content']
        
        if is_excel_source and table_data and '|' not in response_text and format_type == 'text':
            if any(keyword in query_lower for keyword in ['table', 'data', 'row', 'column', 'list', 'show', 'display']):
                columns = table_data['columns']
                data = table_data['data']
                
                table_md = "| " + " | ".join(str(col) for col in columns) + " |\n"
                table_md += "| " + " | ".join(["---"] * len(columns)) + " |\n"
                for row in data:
                    table_md += "| " + " | ".join(str(cell) for cell in row) + " |\n"
                
                response_text += "\n\n" + table_md
        
        return jsonify({
            'response': response_text, 
            'format': format_type,
            'session_id': session_id
        }), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/chat', methods=['POST'])
@login_required
def chat():
    data = request.json
    query = data.get('message', '')
    format_type = data.get('format', 'text')  # text, excel, word, ppt
    session_id = data.get('session_id', str(uuid.uuid4()))  # Use provided session ID or create new
    
    if not query:
        return jsonify({'error': 'No message provided'}), 400
    
    # Get user's accessible companies
    user_companies = session.get('companies', [])
    is_admin = session.get('role') == 'admin'
    
    # Initialize conversation history for new session
    if session_id not in conversation_history:
        conversation_history[session_id] = []
    
    # Add user message to history
    conversation_history[session_id].append({'role': 'user', 'content': query})
    
    # Keep only last 10 messages to manage context window
    if len(conversation_history[session_id]) > 10:
        conversation_history[session_id] = conversation_history[session_id][-10:]
    
    # Load active training rules
    conn = get_db_connection()
    cursor = conn.cursor()
    
    if is_admin:
        cursor.execute('''
            SELECT instruction, category, priority
            FROM ai_training
            WHERE enabled = 1
            ORDER BY priority ASC
        ''')
    else:
        # For non-admin users, load training for their companies + ALL
        placeholders = ','.join(['?'] * len(user_companies))
        cursor.execute(f'''
            SELECT instruction, category, priority
            FROM ai_training
            WHERE enabled = 1 AND (company IN ({placeholders}) OR company = 'ALL')
            ORDER BY priority ASC
        ''', user_companies)
    
    training_rules = cursor.fetchall()
    conn.close()
    
    # Build training instructions for system prompt
    training_instructions = ""
    if training_rules:
        training_instructions = "\n\nADMIN TRAINING RULES (Highest Priority):\n"
        for rule in training_rules:
            training_instructions += f"- [{rule['category']}] {rule['instruction']}\n"
        training_instructions += "\nIMPORTANT: These training rules MUST be followed above all other instructions.\n"
    
    # Get intelligent column mapping for query
    column_mapping = get_column_mapping_for_query(query, user_companies if not is_admin else COMPANIES)
    
    # Build column mapping instructions
    column_mapping_instructions = ""
    if column_mapping:
        column_mapping_instructions = "\n\nCOLUMN MAPPING (Use these actual column names from the data):\n"
        for term, actual_col in column_mapping.items():
            column_mapping_instructions += f"- When user mentions '{term}', use column '{actual_col}'\n"
        column_mapping_instructions += "\nIMPORTANT: Always use the actual column names from the data, not the user's terminology.\n"
    
    # Retrieve relevant documents with company filtering
    query_embedding = embedding_model.encode([query]).tolist()
    
    # For non-admin users, filter by company in the where clause
    if is_admin:
        results = collection.query(
            query_embeddings=query_embedding,
            n_results=10
        )
    else:
        # Filter results by user's accessible companies
        results = collection.query(
            query_embeddings=query_embedding,
            n_results=10,
            where={"company": {"$in": user_companies}}
        )
    
    # Build context from retrieved documents with better organization
    context = ""
    table_data = None
    is_excel_source = False
    source_info = []
    
    if results['documents'] and results['documents'][0]:
        # Organize retrieved documents by source for better context
        for i, (doc, meta) in enumerate(zip(results['documents'][0], results['metadatas'][0])):
            source = meta.get('source', 'Unknown')
            company = meta.get('company', 'Unknown')
            chunk_num = meta.get('chunk', i)
            
            # Add source information
            source_info.append(f"[Source: {source}, Company: {company}, Chunk: {chunk_num}]")
            
            # Check if this is from Excel
            if meta.get('is_excel', False):
                is_excel_source = True
                if meta.get('table_data'):
                    # Parse JSON string back to dictionary
                    try:
                        table_data = json.loads(meta['table_data'])
                    except:
                        table_data = None
            
            # Add document content with source reference
            context += f"\n--- {source_info[-1]} ---\n{doc}\n"
        
        # Add summary of sources
        if source_info:
            context = f"Retrieved from {len(source_info)} document chunks from {len(set(s.split(',')[1].split(':')[1].strip() for s in source_info))} different companies:\n" + context
    
    # Build conversation context from history
    conversation_context = ""
    if len(conversation_history[session_id]) > 1:
        conversation_context = "\n\n".join([
            f"{msg['role'].capitalize()}: {msg['content']}" 
            for msg in conversation_history[session_id][:-1]
        ])
    
    # Create enhanced system prompt with comprehensive formatting and validation requirements
    system_prompt = """You are an intelligent data analysis and presentation assistant. Your job is to provide accurate, well-formatted, professional responses based on user intent.

""" + training_instructions + column_mapping_instructions + """

CRITICAL DATA ACCURACY RULES:

A. STRICT SPORTS MATCHING
- When user asks for specific sport (e.g., "Football", "Cricket", "Basketball"), return ONLY personnel where that sport exists in the Sports field
- Example: "Football players" → Only include records where "Football" is explicitly present in Sports column
- Sports like "Cricket, Football" → Include (Football present)
- Sports like "Football, Basketball" → Include (Football present)
- Sports like "Volleyball, Tennis" → Exclude (Football absent)
- NEVER infer suitability from other sports
- If sport-specific data not available: "[Sport]-specific data is not available. Unable to determine [sport] suitability."

B. QUERY INTENT RULES
- When user asks for sport suitability, evaluate in this order:
  1. Is there a dedicated [Sport] field?
  2. Is there [Sport] Experience field?
  3. Is [Sport] present in Sports Played field?
- If yes to any: Return only matching personnel
- If no sport-related data exists: "[Sport]-specific data is not available. Unable to determine [sport] suitability."
- NEVER recommend based on other sports (Volleyball, Basketball, Tennis, Cricket, etc.)

C. STRICT FILTERING
- When user specifies conditions, EVERY returned record must satisfy ALL conditions
- Do NOT include records that satisfy only one condition
- Example: "Football with service < 5 years" means Football present AND service < 5
- Wrong: Football absent but service < 5, or Football present but service > 5
- Correct: Football present AND service < 5

D. NO ASSUMPTIONS
- NEVER assume suitability without supporting data
- If suitability criteria are not available, state: "Insufficient data available to determine suitability. Showing personnel with [relevant field] instead."
- Never invent suitability scores or infer from unrelated fields

E. HALLUCINATION PREVENTION
- NEVER guess values or create missing records
- NEVER recommend personnel without supporting data
- NEVER infer sports experience from unrelated fields
- If no records match: "No personnel found matching all specified conditions."
- Do not generate approximate results

F. QUERY VERIFICATION
- Before generating response, verify every row matches all filters
- Cross-check that every row satisfies every filter condition
- Only display output after validation passes
- Display "Records Validated: XX, Records Returned: XX" - both numbers must match

G. COLUMN SEPARATION
- Sports and Experience must ALWAYS be separate columns
- Never merge sports and experience into one column
- Correct format: Separate "Sports" column and "Experience (Years)" column

RESPONSE FORMATTING REQUIREMENTS:

A. TABULAR OUTPUT (Use when user asks for: table, tabular report, list of records, comparison, summary report, inventory, employee data, sales data, financial data, any structured dataset)
- Serial Number (S.No.) column MANDATORY as first column
- Clear and meaningful column headings
- One record per row
- Proper alignment and spacing
- Missing values display as "N/A"
- Include totals and summary rows where applicable
- Keep large tables readable and organized

Example table format:
| S.No. | Employee Name | Department | Salary |
| ----- | ------------- | ---------- | ------ |
| 1     | John Doe      | Sales      | 50000  |
| 2     | Jane Smith    | HR         | 60000  |
| 3     | Robert Brown  | IT         | 70000  |

B. NON-TABULAR OUTPUT (Use for: questions, explanations, analysis, summary, recommendations, status report, insights, trends, observations, general information)
- Use clear headings and sub-headings
- Use bullet points and numbering where appropriate
- Group related information together
- Highlight key findings separately
- Keep output professional and easy to read
- Avoid large unstructured paragraphs
- Use sections: Summary, Key Findings, Analysis, Recommendations, Conclusion

Example structured format:
## Summary
- Total employees: 120
- Active employees: 110

## Key Findings
1. Sales department has highest headcount
2. IT department has highest average salary

## Recommendations
1. Increase hiring in Operations
2. Review salary structure for support staff

C. AUTOMATIC FORMAT DETECTION
Determine the best response format based on user intent:
- "Show all employees" → Table
- "List routers with utilization" → Table
- "Compare company performance" → Table + Summary
- "Explain why sales dropped" → Organized text response
- "Give insights from the data" → Organized text response
- "Generate report" → Executive Summary + Tables + Key Findings

D. CONSISTENCY REQUIREMENTS
- Be professional and easy to read
- Avoid raw JSON output
- Avoid dumping unformatted data
- Maintain consistent formatting across all responses
- Never mix multiple formats in a single response unless specifically requested

IMPORTANT SECURITY INSTRUCTIONS:
- If the user asks about data from companies they don't have access to, respond with: "You do not have permission to access that company's data."
- Only provide information from the document context provided
- Do not make up or hallucinate information about companies not in the context
- If the context is empty or doesn't contain relevant information, say so clearly
- Maintain strict data confidentiality and access control
"""

    # Detect format request from user query with automatic format detection
    query_lower = query.lower()
    format_instruction = ""
    
    # Tabular output triggers
    tabular_keywords = ['table', 'tabular report', 'list of records', 'comparison', 'summary report', 'inventory', 'employee data', 'sales data', 'financial data', 'show all', 'list all', 'display all', 'get all', 'show me all', 'what are all', 'who are all']
    
    # Non-tabular output triggers
    non_tabular_keywords = ['explain', 'why', 'how', 'analyze', 'analysis', 'insights', 'trends', 'observations', 'recommendations', 'status report', 'what does', 'what is the meaning', 'tell me about', 'describe', 'give insights', 'generate report']
    
    # Parse filters from query for transparency and validation
    applied_filters = []
    validation_rules = []
    
    # Sports matching - exact matching only
    sports_keywords = ['football', 'basketball', 'cricket', 'volleyball', 'tennis', 'badminton']
    detected_sport = None
    for sport in sports_keywords:
        if sport in query_lower:
            detected_sport = sport.capitalize()
            applied_filters.append(f"Sports contains {detected_sport}")
            validation_rules.append(f"IF {detected_sport} NOT IN Sports THEN Exclude Record")
            break
    
    # Experience/service filters
    if 'experience' in query_lower or 'service' in query_lower:
        if 'less than' in query_lower or '<' in query_lower:
            applied_filters.append("Experience/Service Years < specified value")
            validation_rules.append("IF Experience/Service Years >= specified value THEN Exclude Record")
        elif 'more than' in query_lower or '>' in query_lower:
            applied_filters.append("Experience/Service Years > specified value")
            validation_rules.append("IF Experience/Service Years <= specified value THEN Exclude Record")
    
    # Check for tabular format request
    if any(keyword in query_lower for keyword in tabular_keywords):
        format_instruction = "RESPONSE FORMAT: Tabular. Provide a properly formatted Markdown table with S.No as the first column, clear headings, one record per row, and 'N/A' for missing values. Include totals where applicable. STRICT FILTERING: Ensure every row satisfies ALL specified conditions. COLUMN SEPARATION: Sports and Experience must be separate columns."
    # Check for non-tabular format request
    elif any(keyword in query_lower for keyword in non_tabular_keywords):
        format_instruction = "RESPONSE FORMAT: Non-tabular. Use structured sections: Summary, Key Findings, Analysis, Recommendations, Conclusion. Use headings, bullet points, and numbering. Group related information. STRICT FILTERING: Ensure every data point satisfies ALL specified conditions."
    else:
        # Default to structured format for general queries
        format_instruction = "RESPONSE FORMAT: Structured. Use the most appropriate format based on the content. For structured data, use tables with S.No. For explanations, use organized sections with headings and bullet points. STRICT FILTERING: Ensure every row satisfies ALL specified conditions. COLUMN SEPARATION: Sports and Experience must be separate columns."
    
    # Override with explicit format types if provided
    if format_type == 'excel':
        format_instruction = "RESPONSE FORMAT: JSON. Provide structured JSON data suitable for Excel import. STRICT FILTERING: Ensure every record satisfies ALL specified conditions. COLUMN SEPARATION: Sports and Experience must be separate fields."
    elif format_type == 'word':
        format_instruction = "RESPONSE FORMAT: Document. Use well-structured text with clear headings and paragraphs suitable for Word documents. STRICT FILTERING: Ensure every data point satisfies ALL specified conditions. COLUMN SEPARATION: Sports and Experience must be separate."
    elif format_type == 'ppt':
        format_instruction = "RESPONSE FORMAT: Presentation. Use bullet points and short sections suitable for PowerPoint slides. STRICT FILTERING: Ensure every data point satisfies ALL specified conditions. COLUMN SEPARATION: Sports and Experience must be separate."
    
    # Build filter explanation for transparency
    filter_explanation = ""
    validation_explanation = ""
    if applied_filters:
        filter_explanation = f"\n\nAPPLIED FILTERS:\n" + "\n".join(f"✓ {filter}" for filter in applied_filters)
    if validation_rules:
        validation_explanation = f"\n\nVALIDATION RULES:\n" + "\n".join(f"• {rule}" for rule in validation_rules)
    
    prompt = f"""{system_prompt}

Document Context:
{context}

Question: {query}

{format_instruction}

{filter_explanation}

{validation_explanation}

VALIDATION REQUIREMENT: Before generating the table, verify every row against all filters. Display "Records Validated: XX, Records Returned: XX" where both numbers must match.

Provide a professional, well-formatted response following the specified format requirements. Ensure strict adherence to all filtering conditions and column separation requirements."""
    
    # Build messages for Ollama with conversation history
    messages = [{'role': 'system', 'content': system_prompt}]
    
    # Add conversation history (excluding current question which we'll add separately)
    for msg in conversation_history[session_id][:-1]:
        messages.append({'role': msg['role'], 'content': msg['content']})
    
    # Add current message with full context
    messages.append({'role': 'user', 'content': prompt})
    
    # Get response from Ollama
    try:
        response = ollama.chat(model='llama3', messages=messages)
        response_text = response['message']['content']
        
        # Add assistant response to history
        conversation_history[session_id].append({'role': 'assistant', 'content': response_text})
        
        # If we have table data and the response doesn't already contain a table, format it
        if is_excel_source and table_data and '|' not in response_text and format_type == 'text':
            # Try to detect if the response should include a table
            if any(keyword in query.lower() for keyword in ['table', 'data', 'row', 'column', 'list', 'show', 'display']):
                # Format the table data as Markdown
                columns = table_data['columns']
                data = table_data['data']
                
                # Create Markdown table
                table_md = "| " + " | ".join(str(col) for col in columns) + " |\n"
                table_md += "| " + " | ".join(["---"] * len(columns)) + " |\n"
                for row in data:
                    table_md += "| " + " | ".join(str(cell) for cell in row) + " |\n"
                
                # Append table to response
                response_text += "\n\n**Tabular Data:**\n" + table_md
        
        return jsonify({
            'response': response_text, 
            'format': format_type,
            'session_id': session_id
        }), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/documents', methods=['GET'])
@login_required
def get_documents():
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        
        if session.get('role') == 'admin':
            # Admin can see all files
            cursor.execute('SELECT * FROM files WHERE status = ? ORDER BY upload_date DESC', ('active',))
        else:
            # Regular users can only see files from their companies
            user_companies = session.get('companies', [])
            placeholders = ','.join('?' * len(user_companies))
            cursor.execute(f'SELECT * FROM files WHERE company IN ({placeholders}) AND status = ? ORDER BY upload_date DESC', 
                          user_companies + ['active'])
        
        files = cursor.fetchall()
        conn.close()
        
        documents = []
        for file in files:
            documents.append({
                'id': file['id'],
                'filename': file['filename'],
                'original_filename': file['original_filename'],
                'company': file['company'],
                'file_size': file['file_size'],
                'upload_date': file['upload_date'],
                'uploaded_by': file['uploaded_by']
            })
        
        return jsonify({'documents': documents}), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/clear', methods=['POST'])
@login_required
@admin_required
def clear_documents():
    try:
        chroma_client.delete_collection("documents")
        global collection
        collection = chroma_client.create_collection(name="documents")
        
        # Clear uploads
        for filename in os.listdir(UPLOAD_FOLDER):
            file_path = os.path.join(UPLOAD_FOLDER, filename)
            if os.path.isfile(file_path):
                os.unlink(file_path)
        
        # Clear file records from database
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute('UPDATE files SET status = ? WHERE status = ?', ('deleted', 'active'))
        conn.commit()
        conn.close()
        
        log_audit(session.get('user_id'), 'clear_all', 'document', details='Cleared all documents', ip_address=request.remote_addr)
        
        return jsonify({'message': 'All documents cleared successfully'}), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/clear-conversation', methods=['POST'])
@login_required
def clear_conversation():
    """Clear conversation history for a specific session or all sessions"""
    try:
        data = request.json
        session_id = data.get('session_id')
        
        if session_id:
            # Clear specific session
            if session_id in conversation_history:
                del conversation_history[session_id]
                log_audit(session.get('user_id'), 'clear_conversation', 'chat', session_id, ip_address=request.remote_addr)
                return jsonify({'message': f'Conversation history cleared for session {session_id}'}), 200
            else:
                return jsonify({'error': 'Session not found'}), 404
        else:
            # Clear all conversations
            conversation_history.clear()
            log_audit(session.get('user_id'), 'clear_all_conversations', 'chat', details='Cleared all conversations', ip_address=request.remote_addr)
            return jsonify({'message': 'All conversation histories cleared'}), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/health', methods=['GET'])
def health_check():
    try:
        # Check if Ollama is running
        ollama.list()
        return jsonify({'status': 'healthy', 'ollama': 'connected'}), 200
    except:
        return jsonify({'status': 'unhealthy', 'ollama': 'disconnected'}), 500

# File management APIs (Admin only)
@app.route('/api/admin/files/<int:file_id>', methods=['DELETE'])
@login_required
@admin_required
def delete_file(file_id):
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute('SELECT * FROM files WHERE id = ? AND status = ?', (file_id, 'active'))
        file = cursor.fetchone()
        
        if not file:
            conn.close()
            return jsonify({'error': 'File not found'}), 404
        
        # Delete from ChromaDB
        collection.delete(where={"source": file['filename']})
        
        # Delete physical file
        filepath = os.path.join(UPLOAD_FOLDER, file['filename'])
        if os.path.exists(filepath):
            os.unlink(filepath)
        
        # Mark as deleted in database
        cursor.execute('UPDATE files SET status = ? WHERE id = ?', ('deleted', file_id))
        conn.commit()
        conn.close()
        
        log_audit(session.get('user_id'), 'delete', 'file', str(file_id), 
                  details=f'Deleted file {file["original_filename"]}', ip_address=request.remote_addr)
        
        return jsonify({'message': 'File deleted successfully'}), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/admin/files/<int:file_id>', methods=['PUT'])
@login_required
@admin_required
def replace_file(file_id):
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file part'}), 400
        
        file = request.files['file']
        company = request.form.get('company')
        
        if not company or company not in COMPANIES:
            return jsonify({'error': 'Valid company required'}), 400
        
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute('SELECT * FROM files WHERE id = ? AND status = ?', (file_id, 'active'))
        old_file = cursor.fetchone()
        
        if not old_file:
            conn.close()
            return jsonify({'error': 'File not found'}), 404
        
        # Delete old file from ChromaDB
        collection.delete(where={"source": old_file['filename']})
        
        # Delete old physical file
        old_filepath = os.path.join(UPLOAD_FOLDER, old_file['filename'])
        if os.path.exists(old_filepath):
            os.unlink(old_filepath)
        
        # Upload new file
        new_filename = f"{uuid.uuid4()}_{file.filename}"
        new_filepath = os.path.join(UPLOAD_FOLDER, new_filename)
        file.save(new_filepath)
        
        file_size = os.path.getsize(new_filepath)
        
        # Extract text from new file
        file_ext = new_filename.rsplit('.', 1)[1].lower()
        table_data = None
        if file_ext in ['xlsx', 'xls']:
            text, table_data = extract_text_from_excel(new_filepath)
        elif file_ext == 'pdf':
            text = extract_text_from_pdf(new_filepath)
        elif file_ext == 'docx':
            text = extract_text_from_word(new_filepath)
        
        # Chunk and store in ChromaDB
        chunks = chunk_text(text)
        embeddings = embedding_model.encode(chunks).tolist()
        ids = [f"{new_filename}_{i}" for i in range(len(chunks))]
        metadatas = [{
            "source": new_filename, 
            "chunk": i, 
            "is_excel": file_ext in ['xlsx', 'xls'], 
            "table_data": table_data if i == 0 else None,
            "company": company,
            "uploaded_by": session.get('user_id')
        } for i in range(len(chunks))]
        
        collection.add(
            embeddings=embeddings,
            documents=chunks,
            ids=ids,
            metadatas=metadatas
        )
        
        # Update database
        cursor.execute('''
            UPDATE files SET filename = ?, original_filename = ?, company = ?, 
            file_size = ?, upload_date = CURRENT_TIMESTAMP, uploaded_by = ?
            WHERE id = ?
        ''', (new_filename, file.filename, company, file_size, session.get('user_id'), file_id))
        conn.commit()
        conn.close()
        
        log_audit(session.get('user_id'), 'replace', 'file', str(file_id), 
                  details=f'Replaced {old_file["original_filename"]} with {file.filename} for company {company}', 
                  ip_address=request.remote_addr)
        
        return jsonify({'message': 'File replaced successfully'}), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# User management APIs (Admin only)
@app.route('/api/admin/users', methods=['POST'])
@login_required
@admin_required
def create_user():
    try:
        data = request.json
        username = data.get('username')
        password = data.get('password')
        full_name = data.get('full_name')
        role = data.get('role')
        companies = data.get('companies', [])
        
        if not all([username, password, full_name, role]):
            return jsonify({'error': 'Missing required fields'}), 400
        
        if role not in ['admin', 'user']:
            return jsonify({'error': 'Invalid role'}), 400
        
        if not isinstance(companies, list) or not companies:
            return jsonify({'error': 'Companies must be a non-empty list'}), 400
        
        for company in companies:
            if company not in COMPANIES:
                return jsonify({'error': f'Invalid company: {company}'}), 400
        
        companies_str = ','.join(companies)
        password_hash = hash_password(password)
        
        conn = get_db_connection()
        cursor = conn.cursor()
        
        # Check if username already exists
        cursor.execute('SELECT id FROM users WHERE username = ?', (username,))
        if cursor.fetchone():
            conn.close()
            return jsonify({'error': 'Username already exists'}), 400
        
        cursor.execute('''
            INSERT INTO users (username, password_hash, full_name, role, companies)
            VALUES (?, ?, ?, ?, ?)
        ''', (username, password_hash, full_name, role, companies_str))
        user_id = cursor.lastrowid
        conn.commit()
        conn.close()
        
        log_audit(session.get('user_id'), 'create', 'user', str(user_id), 
                  details=f'Created user {username} with role {role} for companies {companies_str}', 
                  ip_address=request.remote_addr)
        
        return jsonify({
            'message': 'User created successfully',
            'user_id': user_id
        }), 201
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/admin/users', methods=['GET'])
@login_required
@admin_required
def get_users():
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute('SELECT * FROM users ORDER BY created_at DESC')
        users = cursor.fetchall()
        conn.close()
        
        user_list = []
        for user in users:
            user_list.append({
                'id': user['id'],
                'username': user['username'],
                'full_name': user['full_name'],
                'role': user['role'],
                'companies': user['companies'].split(','),
                'status': user['status'],
                'created_at': user['created_at']
            })
        
        return jsonify({'users': user_list}), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/admin/users/<int:user_id>', methods=['PUT'])
@login_required
@admin_required
def update_user(user_id):
    try:
        data = request.json
        full_name = data.get('full_name')
        role = data.get('role')
        companies = data.get('companies', [])
        status = data.get('status')
        
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute('SELECT * FROM users WHERE id = ?', (user_id,))
        user = cursor.fetchone()
        
        if not user:
            conn.close()
            return jsonify({'error': 'User not found'}), 404
        
        # Prevent admin from deactivating themselves
        if user_id == session.get('user_id') and status == 'inactive':
            conn.close()
            return jsonify({'error': 'Cannot deactivate yourself'}), 400
        
        update_fields = []
        update_values = []
        
        if full_name:
            update_fields.append('full_name = ?')
            update_values.append(full_name)
        
        if role:
            if role not in ['admin', 'user']:
                conn.close()
                return jsonify({'error': 'Invalid role'}), 400
            update_fields.append('role = ?')
            update_values.append(role)
        
        if companies:
            if not isinstance(companies, list) or not companies:
                conn.close()
                return jsonify({'error': 'Companies must be a non-empty list'}), 400
            for company in companies:
                if company not in COMPANIES:
                    conn.close()
                    return jsonify({'error': f'Invalid company: {company}'}), 400
            companies_str = ','.join(companies)
            update_fields.append('companies = ?')
            update_values.append(companies_str)
        
        if status:
            if status not in ['active', 'inactive']:
                conn.close()
                return jsonify({'error': 'Invalid status'}), 400
            update_fields.append('status = ?')
            update_values.append(status)
        
        update_fields.append('updated_at = CURRENT_TIMESTAMP')
        update_values.append(user_id)
        
        cursor.execute(f'UPDATE users SET {", ".join(update_fields)} WHERE id = ?', update_values)
        conn.commit()
        conn.close()
        
        log_audit(session.get('user_id'), 'update', 'user', str(user_id), 
                  details=f'Updated user {user["username"]}', ip_address=request.remote_addr)
        
        return jsonify({'message': 'User updated successfully'}), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/admin/users/<int:user_id>', methods=['DELETE'])
@login_required
@admin_required
def delete_user(user_id):
    try:
        # Prevent admin from deleting themselves
        if user_id == session.get('user_id'):
            return jsonify({'error': 'Cannot delete yourself'}), 400
        
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute('SELECT * FROM users WHERE id = ?', (user_id,))
        user = cursor.fetchone()
        
        if not user:
            conn.close()
            return jsonify({'error': 'User not found'}), 404
        
        cursor.execute('DELETE FROM users WHERE id = ?', (user_id,))
        conn.commit()
        conn.close()
        
        log_audit(session.get('user_id'), 'delete', 'user', str(user_id), 
                  details=f'Deleted user {user["username"]}', ip_address=request.remote_addr)
        
        return jsonify({'message': 'User deleted successfully'}), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/admin/users/<int:user_id>/password', methods=['PUT'])
@login_required
@admin_required
def reset_password(user_id):
    try:
        data = request.json
        new_password = data.get('new_password')
        
        if not new_password:
            return jsonify({'error': 'New password required'}), 400
        
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute('SELECT * FROM users WHERE id = ?', (user_id,))
        user = cursor.fetchone()
        
        if not user:
            conn.close()
            return jsonify({'error': 'User not found'}), 404
        
        password_hash = hash_password(new_password)
        cursor.execute('UPDATE users SET password_hash = ?, updated_at = CURRENT_TIMESTAMP WHERE id = ?', 
                      (password_hash, user_id))
        conn.commit()
        conn.close()
        
        log_audit(session.get('user_id'), 'reset_password', 'user', str(user_id), 
                  details=f'Reset password for user {user["username"]}', ip_address=request.remote_addr)
        
        return jsonify({'message': 'Password reset successfully'}), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# Audit log API (Admin only)
@app.route('/api/admin/audit-logs', methods=['GET'])
@login_required
@admin_required
def get_audit_logs():
    try:
        limit = request.args.get('limit', 100, type=int)
        offset = request.args.get('offset', 0, type=int)
        
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute('''
            SELECT al.*, u.username 
            FROM audit_logs al 
            LEFT JOIN users u ON al.user_id = u.id 
            ORDER BY al.timestamp DESC 
            LIMIT ? OFFSET ?
        ''', (limit, offset))
        logs = cursor.fetchall()
        conn.close()
        
        log_list = []
        for log in logs:
            log_list.append({
                'id': log['id'],
                'user_id': log['user_id'],
                'username': log['username'],
                'action': log['action'],
                'resource_type': log['resource_type'],
                'resource_id': log['resource_id'],
                'details': log['details'],
                'ip_address': log['ip_address'],
                'timestamp': log['timestamp']
            })
        
        return jsonify({'logs': log_list}), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/generate/excel', methods=['POST'])
@login_required
def generate_excel():
    """Generate Excel file from AI response"""
    data = request.json
    content = data.get('content', '')
    
    if not content:
        return jsonify({'error': 'No content provided'}), 400
    
    try:
        excel_file = generate_excel_file(content, 'response.xlsx')
        return send_file(
            excel_file,
            mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            as_attachment=True,
            download_name='ai_response.xlsx'
        )
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/generate/word', methods=['POST'])
@login_required
def generate_word():
    """Generate Word file from AI response"""
    data = request.json
    content = data.get('content', '')
    
    if not content:
        return jsonify({'error': 'No content provided'}), 400
    
    try:
        word_file = generate_word_file(content, 'response.docx')
        return send_file(
            word_file,
            mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            as_attachment=True,
            download_name='ai_response.docx'
        )
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# AI Training APIs (Admin only)
@app.route('/api/admin/training', methods=['POST'])
def create_training():
    """Create a new AI training rule"""
    data = request.json
    instruction = data.get('instruction')
    category = data.get('category')
    company = data.get('company', 'ALL')
    priority = data.get('priority', 5)
    
    if not instruction or not category:
        return jsonify({'error': 'Instruction and category are required'}), 400
    
    conn = get_db_connection()
    cursor = conn.cursor()
    
    try:
        # Default to admin user (ID=1) if no session
        user_id = session.get('user_id', 1)
        cursor.execute('''
            INSERT INTO ai_training (company, instruction, category, priority, created_by, enabled)
            VALUES (?, ?, ?, ?, ?, 1)
        ''', (company, instruction, category, priority, user_id))
        
        training_id = cursor.lastrowid
        conn.commit()
        
        log_audit(session.get('user_id'), 'create', 'training', str(training_id), 
                  details=f'Created training rule: {instruction[:50]}...', ip_address=request.remote_addr)
        
        return jsonify({'message': 'Training rule created successfully', 'id': training_id}), 201
    except Exception as e:
        conn.rollback()
        return jsonify({'error': str(e)}), 500
    finally:
        conn.close()

@app.route('/api/admin/training', methods=['GET'])
def get_training():
    """Get all training rules with optional filtering"""
    category = request.args.get('category')
    company = request.args.get('company')
    enabled_only = request.args.get('enabled_only', 'false').lower() == 'true'
    
    conn = get_db_connection()
    cursor = conn.cursor()
    
    try:
        query = '''
            SELECT t.*, u.username as created_by_username
            FROM ai_training t
            LEFT JOIN users u ON t.created_by = u.id
            WHERE 1=1
        '''
        params = []
        
        if category:
            query += ' AND t.category = ?'
            params.append(category)
        
        if company and company != 'ALL':
            query += ' AND (t.company = ? OR t.company = "ALL")'
            params.append(company)
        
        if enabled_only:
            query += ' AND t.enabled = 1'
        
        query += ' ORDER BY t.priority ASC, t.created_date DESC'
        
        cursor.execute(query, params)
        trainings = cursor.fetchall()
        
        result = []
        for t in trainings:
            result.append({
                'id': t['id'],
                'company': t['company'],
                'instruction': t['instruction'],
                'category': t['category'],
                'priority': t['priority'],
                'created_by': t['created_by'],
                'created_by_username': t['created_by_username'],
                'created_date': t['created_date'],
                'updated_date': t['updated_date'],
                'enabled': bool(t['enabled'])
            })
        
        return jsonify(result), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500
    finally:
        conn.close()

@app.route('/api/admin/training/<int:training_id>', methods=['PUT'])
def update_training(training_id):
    """Update an existing training rule"""
    data = request.json
    instruction = data.get('instruction')
    category = data.get('category')
    company = data.get('company')
    priority = data.get('priority')
    enabled = data.get('enabled')
    
    conn = get_db_connection()
    cursor = conn.cursor()
    
    try:
        # Check if training exists
        cursor.execute('SELECT * FROM ai_training WHERE id = ?', (training_id,))
        if not cursor.fetchone():
            return jsonify({'error': 'Training rule not found'}), 404
        
        # Build update query dynamically
        update_fields = []
        params = []
        
        if instruction is not None:
            update_fields.append('instruction = ?')
            params.append(instruction)
        
        if category is not None:
            update_fields.append('category = ?')
            params.append(category)
        
        if company is not None:
            update_fields.append('company = ?')
            params.append(company)
        
        if priority is not None:
            update_fields.append('priority = ?')
            params.append(priority)
        
        if enabled is not None:
            update_fields.append('enabled = ?')
            params.append(1 if enabled else 0)
        
        if update_fields:
            update_fields.append('updated_date = CURRENT_TIMESTAMP')
            params.append(training_id)
            
            cursor.execute(f'''
                UPDATE ai_training
                SET {', '.join(update_fields)}
                WHERE id = ?
            ''', params)
            
            conn.commit()
            
            log_audit(session.get('user_id'), 'update', 'training', str(training_id), 
                      details=f'Updated training rule', ip_address=request.remote_addr)
        
        return jsonify({'message': 'Training rule updated successfully'}), 200
    except Exception as e:
        conn.rollback()
        return jsonify({'error': str(e)}), 500
    finally:
        conn.close()

@app.route('/api/admin/training/<int:training_id>', methods=['DELETE'])
def delete_training(training_id):
    """Delete a training rule"""
    conn = get_db_connection()
    cursor = conn.cursor()
    
    try:
        cursor.execute('DELETE FROM ai_training WHERE id = ?', (training_id,))
        
        if cursor.rowcount == 0:
            return jsonify({'error': 'Training rule not found'}), 404
        
        conn.commit()
        
        log_audit(session.get('user_id'), 'delete', 'training', str(training_id), 
                  details=f'Deleted training rule', ip_address=request.remote_addr)
        
        return jsonify({'message': 'Training rule deleted successfully'}), 200
    except Exception as e:
        conn.rollback()
        return jsonify({'error': str(e)}), 500
    finally:
        conn.close()

@app.route('/api/admin/training/categories', methods=['GET'])
def get_training_categories():
    """Get all unique training categories"""
    conn = get_db_connection()
    cursor = conn.cursor()
    
    try:
        cursor.execute('SELECT DISTINCT category FROM ai_training ORDER BY category')
        categories = [row['category'] for row in cursor.fetchall()]
        
        return jsonify(categories), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500
    finally:
        conn.close()

@app.route('/api/admin/training/load', methods=['GET'])
@login_required
def load_active_training():
    """Load all active training rules for AI processing"""
    user_companies = session.get('companies', [])
    is_admin = session.get('role') == 'admin'
    
    conn = get_db_connection()
    cursor = conn.cursor()
    
    try:
        if is_admin:
            cursor.execute('''
                SELECT instruction, category, priority
                FROM ai_training
                WHERE enabled = 1
                ORDER BY priority ASC
            ''')
        else:
            # For non-admin users, load training for their companies + ALL
            placeholders = ','.join(['?'] * len(user_companies))
            cursor.execute(f'''
                SELECT instruction, category, priority
                FROM ai_training
                WHERE enabled = 1 AND (company IN ({placeholders}) OR company = 'ALL')
                ORDER BY priority ASC
            ''', user_companies)
        
        trainings = cursor.fetchall()
        
        result = []
        for t in trainings:
            result.append({
                'instruction': t['instruction'],
                'category': t['category'],
                'priority': t['priority']
            })
        
        return jsonify(result), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500
    finally:
        conn.close()

@app.route('/api/generate/ppt', methods=['POST'])
@login_required
def generate_ppt():
    """Generate PowerPoint file from AI response"""
    data = request.json
    content = data.get('content', '')
    
    if not content:
        return jsonify({'error': 'No content provided'}), 400
    
    try:
        ppt_file = generate_ppt_file(content, 'response.pptx')
        return send_file(
            ppt_file,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name='ai_response.pptx'
        )
    except Exception as e:
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    print("Starting AI Platform...")
    print("Make sure Ollama is running with: ollama serve")
    print("Make sure Llama 3 is pulled: ollama pull llama3")
    app.run(debug=True, port=5000)
